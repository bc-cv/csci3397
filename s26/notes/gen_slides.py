"""
Generate PowerPoint slides with figures for Lectures 32-36.
Usage: conda run -n csci3397 python gen_slides.py
"""
import io
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
import numpy as np

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── colour palette ──────────────────────────────────────────────────────────
MAROON  = RGBColor(0x80, 0x00, 0x00)
DARK    = RGBColor(0x33, 0x33, 0x33)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF5, 0xF5, 0xF5)
BLUE    = RGBColor(0x2B, 0x6C, 0xB3)
GREEN   = RGBColor(0x27, 0x7A, 0x4D)
ORANGE  = RGBColor(0xD4, 0x7B, 0x2A)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GRAY    = RGBColor(0x99, 0x99, 0x99)

# matplotlib colours
M_MAROON = "#800000"
M_BLUE   = "#2B6CB3"
M_GREEN  = "#277A4D"
M_ORANGE = "#D47B2A"
M_RED    = "#C0392B"
M_GRAY   = "#999999"
M_DARK   = "#333333"
M_LIGHT  = "#F0F0F0"

# ── helpers ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, w, h, text, font_size=18, bold=False,
                 color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, w, h, items, font_size=18,
                     color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def title_slide(prs, lec_num, title, subtitle):
    slide = add_slide(prs)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, MAROON)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1),
                 f"Lecture {lec_num}", 28, False, RGBColor(0xFF,0xCC,0xCC))
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11.7), Inches(2),
                 title, 44, True, WHITE)
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1),
                 subtitle, 22, False, RGBColor(0xFF,0xCC,0xCC))
    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8),
                 "CSCI 3397 / PSYC 3317: Biomedical Image Analysis  |  Spring 2026  |  Donglai Wei",
                 16, False, RGBColor(0xFF,0xCC,0xCC))
    return slide

def section_slide(prs, number, title):
    slide = add_slide(prs)
    add_bg(slide, LIGHT)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), MAROON)
    add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1),
                 f"{number}", 80, True, MAROON)
    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.2),
                 title, 36, True, DARK)
    return slide

def content_slide(prs, title, body_items=None, body_text=None):
    slide = add_slide(prs)
    add_bg(slide, WHITE)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), MAROON)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
                 title, 30, True, MAROON)
    add_rect(slide, Inches(0.6), Inches(1.05), Inches(12), Pt(2), MAROON)
    if body_items:
        add_bullet_slide(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5),
                         body_items, 20, DARK)
    elif body_text:
        add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5),
                     body_text, 20, False, DARK)
    return slide

def code_slide(prs, title, code_text, note=""):
    slide = add_slide(prs)
    add_bg(slide, WHITE)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), MAROON)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
                 title, 30, True, MAROON)
    add_rect(slide, Inches(0.6), Inches(1.05), Inches(12), Pt(2), MAROON)
    # code box
    code_bg = add_rect(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.0),
                       RGBColor(0xF8, 0xF8, 0xF8), GRAY)
    add_text_box(slide, Inches(0.9), Inches(1.5), Inches(11.4), Inches(4.6),
                 code_text, 16, False, DARK, font_name="Courier New")
    if note:
        add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.7),
                     note, 16, False, GRAY)
    return slide

def fig_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf

def figure_slide(prs, title, fig, fig_width=Inches(10), fig_left=Inches(1.6),
                 fig_top=Inches(1.4), note=""):
    slide = add_slide(prs)
    add_bg(slide, WHITE)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), MAROON)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
                 title, 30, True, MAROON)
    add_rect(slide, Inches(0.6), Inches(1.05), Inches(12), Pt(2), MAROON)
    buf = fig_to_image(fig)
    slide.shapes.add_picture(buf, fig_left, fig_top, width=fig_width)
    if note:
        add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
                     note, 14, False, GRAY)
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# FIGURES
# ══════════════════════════════════════════════════════════════════════════════

def fig_agency_spectrum():
    fig, ax = plt.subplots(figsize=(12, 5))
    levels = [
        ("Simple\nProcessor", 0, M_GRAY,   "LLM output has\nno impact on flow"),
        ("Router",            1, M_BLUE,   "LLM controls\nif/else switch"),
        ("Tool\nCall",        2, M_GREEN,  "LLM chooses\nfunction + args"),
        ("Multi-step\nAgent", 3, M_ORANGE, "LLM controls\niteration & loop"),
        ("Multi-\nAgent",     4, M_RED,    "Agent spawns\nother agents"),
        ("Code\nAgent",       5, M_MAROON, "LLM writes code,\ndefines own tools"),
    ]
    stars = ["", "\u2605", "\u2605\u2605", "\u2605\u2605", "\u2605\u2605\u2605", "\u2605\u2605\u2605"]
    ax.set_xlim(-0.8, 6.5)
    ax.set_ylim(-1.5, 3.5)
    # arrow
    ax.annotate("", xy=(5.7, 0), xytext=(-0.5, 0),
                arrowprops=dict(arrowstyle="-|>", lw=2.5, color=M_DARK))
    ax.text(2.5, -1.2, "Increasing Agency  \u2192", ha="center", fontsize=14,
            color=M_DARK, style="italic")
    for label, x, color, desc in levels:
        ax.plot(x, 0, "o", markersize=28, color=color, zorder=5)
        ax.text(x, 0, stars[x], ha="center", va="center", fontsize=10,
                color="white", fontweight="bold", zorder=6)
        ax.text(x, 0.9, label, ha="center", va="bottom", fontsize=12,
                fontweight="bold", color=color)
        ax.text(x, 2.3, desc, ha="center", va="bottom", fontsize=10,
                color=M_DARK, bbox=dict(boxstyle="round,pad=0.3",
                facecolor=color, alpha=0.12, edgecolor=color, linewidth=1))
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_messages_format():
    fig, ax = plt.subplots(figsize=(10, 5.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7)
    msgs = [
        (0.5, 5.8, "system", "You are a radiology assistant.\nUse precise anatomical terms.", M_MAROON, 9.0),
        (0.5, 3.9, "user",   "Summarize this CT report:\n\"2.3cm hypodense lesion in segment VII...\"", M_BLUE, 9.0),
        (0.5, 2.0, "assistant", "The CT shows a 2.3cm hypodense hepatic\nlesion in segment VII, suggestive of...", M_GREEN, 9.0),
        (0.5, 0.5, "user",   "What is the differential diagnosis?", M_BLUE, 9.0),
    ]
    for x, y, role, text, color, w in msgs:
        rect = mpatches.FancyBboxPatch((x, y), w, 1.4, boxstyle="round,pad=0.15",
                                        facecolor=color, alpha=0.10, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + 0.2, y + 1.15, role, fontsize=11, fontweight="bold", color=color)
        ax.text(x + 0.2, y + 0.15, text, fontsize=9.5, color=M_DARK, va="bottom")
    ax.set_title("Messages Format: Conversation as a List", fontsize=14,
                 fontweight="bold", color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_api_call_flow():
    fig, ax = plt.subplots(figsize=(11, 4))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 4)
    boxes = [
        (0.3, 1.2, 2.2, 1.6, "Your Code\n(Python)", M_BLUE),
        (4.0, 1.2, 2.8, 1.6, "Anthropic API\n(cloud)", M_MAROON),
        (8.3, 1.2, 2.4, 1.6, "Claude Model\n(LLM)", M_GREEN),
    ]
    for x, y, w, h, label, color in boxes:
        rect = mpatches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.2",
                                        facecolor=color, alpha=0.15, edgecolor=color, linewidth=2)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h/2, label, ha="center", va="center",
                fontsize=13, fontweight="bold", color=color)
    # arrows
    ax.annotate("", xy=(3.8, 2.3), xytext=(2.7, 2.3),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.text(3.25, 2.6, "messages\n+ params", ha="center", fontsize=9, color=M_DARK)
    ax.annotate("", xy=(2.7, 1.7), xytext=(3.8, 1.7),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.text(3.25, 1.0, "streaming\nresponse", ha="center", fontsize=9, color=M_DARK)
    ax.annotate("", xy=(8.1, 2.3), xytext=(7.0, 2.3),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_GRAY))
    ax.annotate("", xy=(7.0, 1.7), xytext=(8.1, 1.7),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_GRAY))
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_token_context():
    fig, ax = plt.subplots(figsize=(10, 3.5))
    models = ["GPT-4o", "Claude Opus", "Claude Sonnet", "Gemini 1.5", "Llama 3"]
    ctx    = [128, 200, 200, 1000, 128]
    colors = [M_BLUE, M_MAROON, M_ORANGE, M_GREEN, M_GRAY]
    bars = ax.barh(models, ctx, color=colors, height=0.55, alpha=0.8, edgecolor="white", linewidth=1.5)
    for bar, val in zip(bars, ctx):
        ax.text(bar.get_width() + 10, bar.get_y() + bar.get_height()/2,
                f"{val}K", va="center", fontsize=12, fontweight="bold", color=M_DARK)
    ax.set_xlabel("Context Window (thousands of tokens)", fontsize=12, color=M_DARK)
    ax.set_title("Context Window Sizes by Model", fontsize=14, fontweight="bold", color=M_DARK)
    ax.set_xlim(0, 1150)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    return fig

def fig_structured_output():
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 6)
    # left: report
    rect1 = mpatches.FancyBboxPatch((0.3, 0.5), 4.0, 5.0, boxstyle="round,pad=0.2",
                                     facecolor=M_BLUE, alpha=0.08, edgecolor=M_BLUE, linewidth=1.5)
    ax.add_patch(rect1)
    ax.text(2.3, 5.2, "Radiology Report (Input)", ha="center", fontsize=12,
            fontweight="bold", color=M_BLUE)
    report = ("CT Abdomen with contrast:\n"
              "The liver demonstrates a 2.3cm\n"
              "hypodense lesion in segment VII.\n"
              "No biliary dilatation. Spleen is\n"
              "normal in size. Kidneys show no\n"
              "hydronephrosis. Small amount of\n"
              "free fluid in pelvis.")
    ax.text(0.6, 4.7, report, fontsize=9.5, color=M_DARK, va="top", family="monospace")
    # arrow
    ax.annotate("", xy=(5.3, 3.0), xytext=(4.5, 3.0),
                arrowprops=dict(arrowstyle="-|>", lw=2.5, color=M_MAROON))
    ax.text(4.9, 3.5, "LLM", ha="center", fontsize=11, fontweight="bold", color=M_MAROON)
    # right: JSON
    rect2 = mpatches.FancyBboxPatch((5.5, 0.5), 5.2, 5.0, boxstyle="round,pad=0.2",
                                     facecolor=M_GREEN, alpha=0.08, edgecolor=M_GREEN, linewidth=1.5)
    ax.add_patch(rect2)
    ax.text(8.1, 5.2, "Structured Output (JSON)", ha="center", fontsize=12,
            fontweight="bold", color=M_GREEN)
    json_text = ('{"findings": [\n'
                 '  {"structure": "liver",\n'
                 '   "observation": "2.3cm hypodense\n'
                 '     lesion in segment VII",\n'
                 '   "severity": "moderate"},\n'
                 '  {"structure": "spleen",\n'
                 '   "observation": "normal size",\n'
                 '   "severity": "normal"},\n'
                 '  {"structure": "pelvis",\n'
                 '   "observation": "free fluid",\n'
                 '   "severity": "mild"}\n'
                 ']}')
    ax.text(5.8, 4.7, json_text, fontsize=9, color=M_DARK, va="top", family="monospace")
    ax.axis("off")
    fig.tight_layout()
    return fig

# ── Lecture 33 figures ──────────────────────────────────────────────────────

def fig_tool_use_cycle():
    fig, ax = plt.subplots(figsize=(12, 5.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6.5)
    # step boxes
    steps = [
        (0.3, 4.0, 3.0, 2.0, "1. Send Request", "messages + tool\nschemas \u2192 LLM", M_BLUE),
        (4.3, 4.0, 3.0, 2.0, "2. LLM Responds", "tool_call: count_cells\n{image: \"slide.png\"}", M_MAROON),
        (8.3, 4.0, 3.0, 2.0, "3. Execute Tool", "Your code runs\ncount_cells(\"slide.png\")", M_GREEN),
        (4.3, 0.8, 3.0, 2.0, "4. Send Result", "tool_result:\n\"Found 247 cells\"", M_ORANGE),
    ]
    for x, y, w, h, title, desc, color in steps:
        rect = mpatches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.2",
                                        facecolor=color, alpha=0.12, edgecolor=color, linewidth=2)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h - 0.3, title, ha="center", va="top",
                fontsize=12, fontweight="bold", color=color)
        ax.text(x + w/2, y + 0.4, desc, ha="center", va="bottom",
                fontsize=10, color=M_DARK, family="monospace")
    # arrows
    ax.annotate("", xy=(4.1, 5.0), xytext=(3.5, 5.0),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.annotate("", xy=(8.1, 5.0), xytext=(7.5, 5.0),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.annotate("", xy=(8.1, 3.8), xytext=(8.1, 2.8),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    # curved arrow back to step 1
    ax.annotate("", xy=(1.8, 3.8), xytext=(4.1, 1.8),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK,
                               connectionstyle="arc3,rad=0.3"))
    ax.text(1.8, 2.8, "Loop until\ndone", ha="center", fontsize=11,
            fontweight="bold", color=M_DARK, style="italic")
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_react_loop():
    fig, ax = plt.subplots(figsize=(11, 5.5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 7)
    # Central cycle
    cx, cy, r = 5.5, 3.5, 2.0
    angles = [90, 210, 330]
    labels = ["Think", "Act", "Observe"]
    colors = [M_MAROON, M_BLUE, M_GREEN]
    descs  = ["Reason about\nwhat to do next", "Call a tool\nwith arguments", "See the\ntool result"]
    coords = []
    for ang, label, color, desc in zip(angles, labels, colors, descs):
        rad = np.radians(ang)
        x = cx + r * np.cos(rad)
        y = cy + r * np.sin(rad)
        coords.append((x, y))
        circle = plt.Circle((x, y), 0.7, color=color, alpha=0.15, linewidth=2, edgecolor=color)
        ax.add_patch(circle)
        ax.text(x, y + 0.1, label, ha="center", va="center", fontsize=14,
                fontweight="bold", color=color)
        # desc outside
        dx = cx + (r + 1.5) * np.cos(rad)
        dy = cy + (r + 1.5) * np.sin(rad)
        ax.text(dx, dy, desc, ha="center", va="center", fontsize=10, color=M_DARK,
                bbox=dict(boxstyle="round,pad=0.3", facecolor=color, alpha=0.06,
                         edgecolor="none"))
    # arrows between nodes
    for i in range(3):
        x1, y1 = coords[i]
        x2, y2 = coords[(i + 1) % 3]
        mx, my = (x1 + x2) / 2, (y1 + y2) / 2
        # offset toward center slightly
        mx = mx * 0.95 + cx * 0.05
        my = my * 0.95 + cy * 0.05
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK,
                                   connectionstyle="arc3,rad=0.35"))
    ax.text(cx, cy, "ReAct\nLoop", ha="center", va="center", fontsize=13,
            fontweight="bold", color=M_DARK, style="italic")
    ax.set_title("ReAct: Reasoning + Acting", fontsize=14, fontweight="bold",
                 color=M_DARK, pad=10)
    ax.axis("off")
    ax.set_aspect("equal")
    fig.tight_layout()
    return fig

def fig_tool_registry():
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 6)
    # registry box
    rect = mpatches.FancyBboxPatch((3.5, 0.5), 4.0, 5.0, boxstyle="round,pad=0.3",
                                    facecolor=M_MAROON, alpha=0.08, edgecolor=M_MAROON, linewidth=2)
    ax.add_patch(rect)
    ax.text(5.5, 5.2, "Tool Registry", ha="center", fontsize=14, fontweight="bold", color=M_MAROON)
    tools_list = ["register_tool()", "get_tool(name)", "get_tool_schemas()", "execute_tool(name, params)"]
    for i, t in enumerate(tools_list):
        ax.text(5.5, 4.3 - i * 0.7, t, ha="center", fontsize=11, color=M_DARK, family="monospace")
    # left: tools registering
    left_tools = [("Read", M_GREEN, 0.5, 4.5), ("Write", M_BLUE, 0.5, 3.5),
                  ("Bash", M_ORANGE, 0.5, 2.5), ("Grep", M_GREEN, 0.5, 1.5)]
    for name, color, x, y in left_tools:
        rect = mpatches.FancyBboxPatch((x, y - 0.3), 1.8, 0.6, boxstyle="round,pad=0.1",
                                        facecolor=color, alpha=0.15, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + 0.9, y, name, ha="center", va="center", fontsize=11,
                fontweight="bold", color=color)
        ax.annotate("", xy=(3.3, y), xytext=(2.5, y),
                    arrowprops=dict(arrowstyle="-|>", lw=1.5, color=M_GRAY))
    # right: agent consuming
    rect = mpatches.FancyBboxPatch((8.5, 2.0), 2.0, 2.0, boxstyle="round,pad=0.2",
                                    facecolor=M_BLUE, alpha=0.12, edgecolor=M_BLUE, linewidth=2)
    ax.add_patch(rect)
    ax.text(9.5, 3.0, "Agent\nLoop", ha="center", va="center", fontsize=13,
            fontweight="bold", color=M_BLUE)
    ax.annotate("", xy=(8.3, 3.3), xytext=(7.7, 3.3),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.text(8.0, 3.7, "schemas", ha="center", fontsize=9, color=M_DARK)
    ax.annotate("", xy=(7.7, 2.7), xytext=(8.3, 2.7),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.text(8.0, 2.3, "dispatch", ha="center", fontsize=9, color=M_DARK)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_code_vs_json():
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))
    for ax in (ax1, ax2):
        ax.set_xlim(0, 6)
        ax.set_ylim(0, 6)
        ax.axis("off")
    # JSON side
    rect1 = mpatches.FancyBboxPatch((0.2, 0.2), 5.6, 5.6, boxstyle="round,pad=0.3",
                                     facecolor=M_BLUE, alpha=0.06, edgecolor=M_BLUE, linewidth=2)
    ax1.add_patch(rect1)
    ax1.set_title("JSON Tool Calling", fontsize=14, fontweight="bold", color=M_BLUE)
    json_code = ('{"tool": "count_cells",\n'
                 ' "args": {"path": "s1.png"}}\n'
                 '\n'
                 '{"tool": "count_cells",\n'
                 ' "args": {"path": "s2.png"}}\n'
                 '\n'
                 '{"tool": "compute_stats",\n'
                 ' "args": {"values":\n'
                 '   [result1, result2]}}')
    ax1.text(0.5, 5.0, json_code, fontsize=10, color=M_DARK, va="top", family="monospace")
    ax1.text(3.0, 0.5, "\u2717 Can't compose   \u2717 No variables",
             ha="center", fontsize=10, color=M_RED, fontweight="bold")
    # Code side
    rect2 = mpatches.FancyBboxPatch((0.2, 0.2), 5.6, 5.6, boxstyle="round,pad=0.3",
                                     facecolor=M_GREEN, alpha=0.06, edgecolor=M_GREEN, linewidth=2)
    ax2.add_patch(rect2)
    ax2.set_title("Code Agent", fontsize=14, fontweight="bold", color=M_GREEN)
    py_code = ('results = []\n'
               'for slide in glob("*.png"):\n'
               '    n = count_cells(slide)\n'
               '    results.append(n)\n'
               '\n'
               'stats = compute_stats(results)\n'
               'print(f"Mean: {stats.mean}")\n'
               '\n'
               '# Composable & natural!')
    ax2.text(0.5, 5.0, py_code, fontsize=10, color=M_DARK, va="top", family="monospace")
    ax2.text(3.0, 0.5, "\u2713 Composable   \u2713 Variables   \u2713 Loops",
             ha="center", fontsize=10, color=M_GREEN, fontweight="bold")
    fig.tight_layout()
    return fig

# ── Lecture 34 figures ──────────────────────────────────────────────────────

def fig_tool_suite():
    fig, ax = plt.subplots(figsize=(12, 5.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6.5)
    categories = [
        ("File I/O", 1.5, [("Read", M_GREEN), ("Write", M_BLUE), ("Edit", M_ORANGE), ("Glob", M_GREEN)]),
        ("Search", 5.0, [("Grep", M_GREEN), ("WebSearch", M_GREEN), ("WebFetch", M_GREEN)]),
        ("Execution", 8.5, [("Bash", M_RED)]),
        ("Notebook", 11.0, [("NotebookEdit", M_ORANGE)]),
    ]
    for cat_name, cx, tools in categories:
        ax.text(cx, 6.0, cat_name, ha="center", fontsize=13, fontweight="bold", color=M_DARK)
        for i, (name, color) in enumerate(tools):
            y = 5.0 - i * 1.1
            rect = mpatches.FancyBboxPatch((cx - 1.0, y - 0.35), 2.0, 0.7,
                                            boxstyle="round,pad=0.1",
                                            facecolor=color, alpha=0.15,
                                            edgecolor=color, linewidth=1.5)
            ax.add_patch(rect)
            ax.text(cx, y, name, ha="center", va="center", fontsize=11,
                    fontweight="bold", color=color)
    # legend
    ax.text(1.5, 0.3, "\u25cf read_only (auto-approve)", fontsize=10, color=M_GREEN)
    ax.text(5.0, 0.3, "\u25cf write (needs permission)", fontsize=10, color=M_BLUE)
    ax.text(8.5, 0.3, "\u25cf dangerous (always ask)", fontsize=10, color=M_RED)
    ax.set_title("The Tool Suite", fontsize=14, fontweight="bold", color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_permission_model():
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 6)
    modes = [
        (1.5, "manual", "Ask for\neverything", M_RED, "X"),
        (5.5, "auto", "Auto-approve reads\nAsk for writes", M_ORANGE, "*"),
        (9.5, "accept-all", "Approve\neverything", M_GREEN, "!"),
    ]
    for x, name, desc, color, icon in modes:
        rect = mpatches.FancyBboxPatch((x - 1.5, 1.5), 3.0, 3.5, boxstyle="round,pad=0.2",
                                        facecolor=color, alpha=0.10, edgecolor=color, linewidth=2)
        ax.add_patch(rect)
        ax.text(x, 4.5, icon, ha="center", fontsize=24)
        ax.text(x, 3.7, name, ha="center", fontsize=14, fontweight="bold", color=color)
        ax.text(x, 2.5, desc, ha="center", fontsize=11, color=M_DARK)
    # safety arrow
    ax.annotate("", xy=(9.5, 0.8), xytext=(1.5, 0.8),
                arrowprops=dict(arrowstyle="-|>", lw=2.5, color=M_DARK))
    ax.text(5.5, 0.4, "More Convenience  \u2192\n\u2190  More Safety", ha="center",
            fontsize=11, color=M_DARK, fontweight="bold")
    ax.set_title("Permission Modes", fontsize=14, fontweight="bold", color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_system_prompt():
    fig, ax = plt.subplots(figsize=(11, 5.5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 7)
    # central box
    rect = mpatches.FancyBboxPatch((3.5, 1.0), 4.0, 5.0, boxstyle="round,pad=0.3",
                                    facecolor=M_MAROON, alpha=0.10, edgecolor=M_MAROON, linewidth=2)
    ax.add_patch(rect)
    ax.text(5.5, 5.7, "System Prompt", ha="center", fontsize=14, fontweight="bold", color=M_MAROON)
    parts = ["Identity & role", "Date & environment", "Git context",
             "CLAUDE.md rules", "Memory index", "Tool instructions"]
    for i, p in enumerate(parts):
        ax.text(5.5, 5.0 - i * 0.7, f"\u2022 {p}", ha="center", fontsize=11, color=M_DARK)
    # sources
    sources = [
        (0.8, 5.5, "config.py", M_BLUE),
        (0.8, 4.0, "git status\ngit log", M_GREEN),
        (0.8, 2.5, "CLAUDE.md", M_ORANGE),
        (10.0, 5.5, "MEMORY.md", M_BLUE),
        (10.0, 4.0, "tool_registry\n.get_schemas()", M_GREEN),
        (10.0, 2.5, "os.environ\nplatform", M_GRAY),
    ]
    for x, y, label, color in sources:
        rect = mpatches.FancyBboxPatch((x - 0.9, y - 0.4), 1.8, 0.8,
                                        boxstyle="round,pad=0.1",
                                        facecolor=color, alpha=0.12,
                                        edgecolor=color, linewidth=1)
        ax.add_patch(rect)
        ax.text(x, y, label, ha="center", va="center", fontsize=9,
                fontweight="bold", color=color)
        target_x = 3.3 if x < 5 else 7.7
        ax.annotate("", xy=(target_x, 3.5), xytext=(x + (0.9 if x < 5 else -0.9), y),
                    arrowprops=dict(arrowstyle="-|>", lw=1.2, color=M_GRAY,
                                   connectionstyle="arc3,rad=0.2"))
    ax.text(5.5, 0.4, "context.py (165 lines) \u2014 assembles all sources into one prompt",
            ha="center", fontsize=11, color=M_DARK, style="italic")
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_compaction():
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6)
    # context window bar
    ax.barh(4.5, 10, left=0.5, height=0.8, color=M_LIGHT, edgecolor=M_DARK, linewidth=1.5)
    ax.barh(4.5, 7, left=0.5, height=0.8, color=M_MAROON, alpha=0.3, edgecolor="none")
    ax.text(4.0, 4.5, "70% used", ha="center", va="center", fontsize=11, color=M_MAROON, fontweight="bold")
    ax.text(9.0, 4.5, "30% free", ha="center", va="center", fontsize=11, color=M_DARK)
    ax.plot([7.85, 7.85], [4.0, 5.0], '--', color=M_RED, lw=2)
    ax.text(7.85, 3.7, "70% threshold", ha="center", fontsize=9, color=M_RED)
    # Layer 1
    rect1 = mpatches.FancyBboxPatch((0.5, 1.5), 4.5, 1.5, boxstyle="round,pad=0.2",
                                     facecolor=M_BLUE, alpha=0.10, edgecolor=M_BLUE, linewidth=1.5)
    ax.add_patch(rect1)
    ax.text(2.75, 2.7, "Layer 1: Snip (free)", ha="center", fontsize=12,
            fontweight="bold", color=M_BLUE)
    ax.text(2.75, 1.9, "Truncate old tool results:\nkeep first 50% + last 25%",
            ha="center", fontsize=10, color=M_DARK)
    # Layer 2
    rect2 = mpatches.FancyBboxPatch((6.0, 1.5), 5.0, 1.5, boxstyle="round,pad=0.2",
                                     facecolor=M_ORANGE, alpha=0.10, edgecolor=M_ORANGE, linewidth=1.5)
    ax.add_patch(rect2)
    ax.text(8.5, 2.7, "Layer 2: AI Summarize (1 API call)", ha="center", fontsize=12,
            fontweight="bold", color=M_ORANGE)
    ax.text(8.5, 1.9, "Split 70/30, summarize old portion\nReplace with compact summary",
            ha="center", fontsize=10, color=M_DARK)
    # arrow
    ax.annotate("", xy=(5.8, 2.25), xytext=(5.2, 2.25),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.text(5.5, 2.7, "still\ntoo big?", ha="center", fontsize=9, color=M_DARK)
    ax.set_title("Two-Layer Context Compaction", fontsize=14, fontweight="bold", color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

# ── Lecture 35 figures ──────────────────────────────────────────────────────

def fig_rag_pipeline():
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 5.5)
    steps = [
        (0.5,  2.0, "Documents\n(papers, reports)", M_GRAY),
        (2.8,  2.0, "Chunk\n(split text)", M_BLUE),
        (5.1,  2.0, "Embed\n(vectors)", M_GREEN),
        (7.4,  2.0, "Index\n(vector DB)", M_ORANGE),
        (9.7,  2.0, "Retrieve +\nGenerate", M_MAROON),
    ]
    for x, y, label, color in steps:
        rect = mpatches.FancyBboxPatch((x, y - 0.6), 2.0, 1.8, boxstyle="round,pad=0.15",
                                        facecolor=color, alpha=0.12, edgecolor=color, linewidth=2)
        ax.add_patch(rect)
        ax.text(x + 1.0, y + 0.3, label, ha="center", va="center",
                fontsize=11, fontweight="bold", color=color)
    for i in range(4):
        x1 = steps[i][0] + 2.1
        x2 = steps[i+1][0] - 0.1
        ax.annotate("", xy=(x2, 2.0), xytext=(x1, 2.0),
                    arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    # query arrow
    ax.annotate("", xy=(9.7, 4.5), xytext=(5.5, 4.5),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_BLUE, linestyle="--"))
    ax.text(7.6, 4.8, "Query embedding", ha="center", fontsize=10, color=M_BLUE,
            fontweight="bold")
    rect_q = mpatches.FancyBboxPatch((3.5, 4.0), 2.0, 1.0, boxstyle="round,pad=0.15",
                                      facecolor=M_BLUE, alpha=0.12, edgecolor=M_BLUE, linewidth=1.5)
    ax.add_patch(rect_q)
    ax.text(4.5, 4.5, "User Query", ha="center", va="center", fontsize=11,
            fontweight="bold", color=M_BLUE)
    # output
    ax.text(10.7, 0.8, "Answer grounded\nin retrieved docs", ha="center",
            fontsize=10, color=M_MAROON, fontweight="bold",
            bbox=dict(boxstyle="round,pad=0.3", facecolor=M_MAROON, alpha=0.08, edgecolor=M_MAROON))
    ax.annotate("", xy=(10.7, 1.2), xytext=(10.7, 1.8),
                arrowprops=dict(arrowstyle="-|>", lw=1.5, color=M_MAROON))
    ax.set_title("RAG Pipeline: Retrieval-Augmented Generation", fontsize=14,
                 fontweight="bold", color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_memory_types():
    fig, ax = plt.subplots(figsize=(11, 5.5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 6.5)
    types = [
        (1.5, 5.0, "user",      "Role, preferences,\nexpertise level", M_BLUE),
        (4.5, 5.0, "feedback",  "Corrections &\nconfirmed approaches", M_GREEN),
        (7.5, 5.0, "project",   "Ongoing work,\ndecisions, deadlines", M_ORANGE),
        (10.0, 5.0, "reference", "External system\npointers", M_MAROON),
    ]
    for x, y, name, desc, color in types:
        rect = mpatches.FancyBboxPatch((x - 1.2, y - 1.5), 2.4, 2.2, boxstyle="round,pad=0.2",
                                        facecolor=color, alpha=0.10, edgecolor=color, linewidth=2)
        ax.add_patch(rect)
        ax.text(x, y, name, ha="center", fontsize=13, fontweight="bold", color=color)
        ax.text(x, y - 0.8, desc, ha="center", fontsize=9.5, color=M_DARK)
    # storage
    rect = mpatches.FancyBboxPatch((2.5, 0.3), 6.0, 1.5, boxstyle="round,pad=0.2",
                                    facecolor=M_DARK, alpha=0.06, edgecolor=M_DARK, linewidth=1.5)
    ax.add_patch(rect)
    ax.text(5.5, 1.3, "MEMORY.md (index) + individual .md files (content)", ha="center",
            fontsize=11, fontweight="bold", color=M_DARK)
    ax.text(5.5, 0.7, "User-scope (~/.nano_claude/memory/)  |  Project-scope (.nano_claude/memory/)",
            ha="center", fontsize=9, color=M_GRAY)
    for x, _, _, _, _ in types:
        ax.annotate("", xy=(5.5, 1.9), xytext=(x, 3.3),
                    arrowprops=dict(arrowstyle="-|>", lw=1, color=M_GRAY,
                                   connectionstyle="arc3,rad=0.15"))
    ax.set_title("Memory Types", fontsize=14, fontweight="bold", color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_multi_agent():
    fig, ax = plt.subplots(figsize=(12, 6))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 7)
    # main agent
    rect = mpatches.FancyBboxPatch((4.0, 5.0), 4.0, 1.5, boxstyle="round,pad=0.2",
                                    facecolor=M_MAROON, alpha=0.15, edgecolor=M_MAROON, linewidth=2)
    ax.add_patch(rect)
    ax.text(6.0, 5.75, "Main Agent", ha="center", va="center",
            fontsize=14, fontweight="bold", color=M_MAROON)
    # sub-agents
    subs = [
        (0.5, 2.5, "Triage\nAgent", M_BLUE, "Scan slides,\nflag abnormal"),
        (3.5, 2.5, "Analysis\nAgent", M_GREEN, "Cell counting,\nmorphometry"),
        (6.5, 2.5, "Literature\nAgent", M_ORANGE, "PubMed search,\nRAG retrieval"),
        (9.5, 2.5, "Report\nAgent", M_MAROON, "Compile\ndiagnostic report"),
    ]
    for x, y, name, color, desc in subs:
        rect = mpatches.FancyBboxPatch((x, y - 0.6), 2.2, 1.8, boxstyle="round,pad=0.15",
                                        facecolor=color, alpha=0.12, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + 1.1, y + 0.5, name, ha="center", va="center",
                fontsize=11, fontweight="bold", color=color)
        ax.text(x + 1.1, y - 0.3, desc, ha="center", va="center", fontsize=8.5, color=M_DARK)
        ax.annotate("", xy=(x + 1.1, y + 1.2), xytext=(6.0, 4.8),
                    arrowprops=dict(arrowstyle="-|>", lw=1.5, color=M_DARK,
                                   connectionstyle="arc3,rad=0.1"))
    # parallel bracket
    ax.annotate("", xy=(8.7, 1.7), xytext=(0.5, 1.7),
                arrowprops=dict(arrowstyle="<->", lw=1.5, color=M_BLUE, linestyle="--"))
    ax.text(4.6, 1.3, "run in parallel (ThreadPoolExecutor)", ha="center",
            fontsize=10, color=M_BLUE, fontweight="bold")
    # report waits
    ax.annotate("", xy=(10.6, 1.7), xytext=(9.2, 1.7),
                arrowprops=dict(arrowstyle="<-", lw=1.5, color=M_RED, linestyle="--"))
    ax.text(10.6, 1.3, "waits for\nall results", ha="center", fontsize=9, color=M_RED)
    ax.set_title("Multi-Agent Pathology Pipeline", fontsize=14, fontweight="bold",
                 color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_mcp_protocol():
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 6)
    # agent
    rect = mpatches.FancyBboxPatch((0.5, 2.0), 2.5, 2.0, boxstyle="round,pad=0.2",
                                    facecolor=M_BLUE, alpha=0.12, edgecolor=M_BLUE, linewidth=2)
    ax.add_patch(rect)
    ax.text(1.75, 3.0, "Agent\n(MCP Client)", ha="center", va="center",
            fontsize=12, fontweight="bold", color=M_BLUE)
    # MCP protocol
    ax.annotate("", xy=(4.2, 3.4), xytext=(3.2, 3.4),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.annotate("", xy=(3.2, 2.6), xytext=(4.2, 2.6),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.text(3.7, 3.8, "JSON-RPC", ha="center", fontsize=9, color=M_DARK, fontweight="bold")
    ax.text(3.7, 2.2, "results", ha="center", fontsize=9, color=M_DARK)
    # MCP servers
    servers = [
        (5.5, 4.0, "PACS Server", "retrieve_study()\nget_dicom()", M_MAROON),
        (5.5, 2.0, "EHR Server", "get_patient()\nget_labs()", M_GREEN),
        (8.5, 4.0, "Lab System", "get_stain_protocol()\nsubmit_order()", M_ORANGE),
        (8.5, 2.0, "Literature", "search_pubmed()\nget_abstract()", M_BLUE),
    ]
    for x, y, name, tools, color in servers:
        rect = mpatches.FancyBboxPatch((x - 1.2, y - 0.7), 2.4, 1.4, boxstyle="round,pad=0.15",
                                        facecolor=color, alpha=0.10, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y + 0.2, name, ha="center", fontsize=10, fontweight="bold", color=color)
        ax.text(x, y - 0.35, tools, ha="center", fontsize=7.5, color=M_DARK, family="monospace")
    # arrows from protocol to servers
    for x, y, _, _, _ in servers:
        sx = 4.5
        ax.annotate("", xy=(x - 1.2, y), xytext=(sx, 3.0),
                    arrowprops=dict(arrowstyle="-", lw=1, color=M_GRAY,
                                   connectionstyle="arc3,rad=0.1"))
    ax.set_title("MCP: Model Context Protocol", fontsize=14, fontweight="bold",
                 color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

# ── Lecture 36 figures ──────────────────────────────────────────────────────

def fig_full_architecture():
    fig, ax = plt.subplots(figsize=(13, 7))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 8)
    # REPL layer
    rect = mpatches.FancyBboxPatch((1.0, 6.5), 11.0, 1.2, boxstyle="round,pad=0.2",
                                    facecolor=M_MAROON, alpha=0.12, edgecolor=M_MAROON, linewidth=2)
    ax.add_patch(rect)
    ax.text(6.5, 7.1, "REPL / UI  (nano_claude.py \u2014 1,985 lines)", ha="center",
            fontsize=13, fontweight="bold", color=M_MAROON)
    # Agent loop
    rect = mpatches.FancyBboxPatch((3.0, 4.8), 7.0, 1.2, boxstyle="round,pad=0.2",
                                    facecolor=M_BLUE, alpha=0.12, edgecolor=M_BLUE, linewidth=2)
    ax.add_patch(rect)
    ax.text(6.5, 5.4, "Core Agent Loop  (agent.py \u2014 179 lines)", ha="center",
            fontsize=13, fontweight="bold", color=M_BLUE)
    ax.annotate("", xy=(6.5, 6.3), xytext=(6.5, 6.1),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    # middle layer
    mid = [
        (1.5, 3.3, "Providers\n618 L", M_GREEN),
        (4.0, 3.3, "Tool Registry\n98 L", M_ORANGE),
        (6.5, 3.3, "Context\n165 L", M_BLUE),
        (9.0, 3.3, "Compaction\n196 L", M_MAROON),
        (11.5, 3.3, "Tools\n1,064 L", M_GREEN),
    ]
    for x, y, label, color in mid:
        rect = mpatches.FancyBboxPatch((x - 1.0, y - 0.5), 2.0, 1.0, boxstyle="round,pad=0.1",
                                        facecolor=color, alpha=0.12, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha="center", va="center", fontsize=9.5,
                fontweight="bold", color=color)
        ax.annotate("", xy=(x, y + 0.6), xytext=(x, 4.6),
                    arrowprops=dict(arrowstyle="-", lw=1, color=M_GRAY))
    # extensions
    ext = [
        (1.5, 1.2, "Memory\n5 modules", M_BLUE),
        (4.0, 1.2, "Multi-Agent\n2 modules", M_ORANGE),
        (6.5, 1.2, "MCP\n4 modules", M_GREEN),
        (9.0, 1.2, "Skills\n4 modules", M_MAROON),
        (11.5, 1.2, "Plugins\n4 modules", M_GRAY),
    ]
    for x, y, label, color in ext:
        rect = mpatches.FancyBboxPatch((x - 1.0, y - 0.5), 2.0, 1.0, boxstyle="round,pad=0.1",
                                        facecolor=color, alpha=0.12, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha="center", va="center", fontsize=9.5,
                fontweight="bold", color=color)
    # bracket label
    ax.text(6.5, 0.2, "Extensions (~3,000 lines)", ha="center", fontsize=11,
            color=M_DARK, fontweight="bold")
    ax.text(6.5, 2.6, "Core (~2,100 lines)", ha="center", fontsize=11,
            color=M_DARK, fontweight="bold")
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_agent_landscape():
    fig, ax = plt.subplots(figsize=(12, 6))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 7)
    # axes labels
    ax.annotate("", xy=(11.5, 0.5), xytext=(0.5, 0.5),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.annotate("", xy=(0.5, 6.5), xytext=(0.5, 0.5),
                arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.text(6.0, 0.1, "Flexibility / Agency  \u2192", ha="center", fontsize=12,
            color=M_DARK, fontweight="bold")
    ax.text(0.1, 3.5, "Complexity\n\u2191", ha="center", va="center", fontsize=12,
            color=M_DARK, fontweight="bold", rotation=0)
    agents = [
        (2.0, 1.5, "LangChain\nChains", M_GRAY, 12),
        (3.5, 2.5, "smolagents", M_GREEN, 13),
        (5.0, 2.0, "OpenAI\nAgents SDK", M_BLUE, 12),
        (7.0, 3.5, "nano-claude\n-code", M_ORANGE, 14),
        (8.5, 5.0, "Claude\nCode", M_MAROON, 14),
        (9.5, 4.5, "Cursor", M_BLUE, 13),
        (10.5, 5.5, "Devin", M_RED, 13),
        (4.5, 4.0, "LangGraph", M_GRAY, 12),
        (6.0, 5.0, "CrewAI", M_ORANGE, 12),
    ]
    for x, y, label, color, fs in agents:
        circle = plt.Circle((x, y), 0.6, color=color, alpha=0.12, linewidth=1.5, edgecolor=color)
        ax.add_patch(circle)
        ax.text(x, y, label, ha="center", va="center", fontsize=fs - 4,
                fontweight="bold", color=color)
    # highlight nano-claude-code
    circle = plt.Circle((7.0, 3.5), 0.7, color=M_ORANGE, alpha=0.0, linewidth=3,
                         edgecolor=M_ORANGE, linestyle="--")
    ax.add_patch(circle)
    ax.text(7.0, 2.5, "We built this!", ha="center", fontsize=11, color=M_ORANGE,
            fontweight="bold", style="italic")
    ax.set_title("Agent Landscape (2025-2026)", fontsize=14, fontweight="bold",
                 color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_deployment_stack():
    fig, ax = plt.subplots(figsize=(11, 6))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 7)
    layers = [
        (0.5, 5.5, 10.0, 1.0, "User Interface", "Web app, CLI, EHR integration", M_BLUE),
        (0.5, 4.2, 10.0, 1.0, "Agent Orchestration", "Agent loop, tool dispatch, permissions", M_MAROON),
        (0.5, 2.9, 10.0, 1.0, "LLM Provider", "Claude API / Azure OpenAI / Local Ollama", M_GREEN),
        (0.5, 1.6, 10.0, 1.0, "Tools & Data", "PACS, EHR, lab systems (via MCP)", M_ORANGE),
        (0.5, 0.3, 10.0, 1.0, "Infrastructure", "Cloud / on-prem, GPU, storage, networking", M_GRAY),
    ]
    for x, y, w, h, name, desc, color in layers:
        rect = mpatches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.15",
                                        facecolor=color, alpha=0.12, edgecolor=color, linewidth=2)
        ax.add_patch(rect)
        ax.text(x + 0.3, y + h/2 + 0.15, name, va="center", fontsize=13,
                fontweight="bold", color=color)
        ax.text(x + 0.3, y + h/2 - 0.25, desc, va="center", fontsize=10, color=M_DARK)
    # right side: concerns
    concerns = [
        (5.5, "\u26d4 HIPAA: PHI in API calls?"),
        (4.2, "\u2696 FDA: SaMD classification?"),
        (2.9, "\u231a Latency: multi-step round trips"),
        (1.6, "\U0001f512 Auth: credential management"),
        (0.3, "\U0001f4ca Monitoring: drift detection"),
    ]
    ax.set_title("Deployment Stack & Concerns", fontsize=14, fontweight="bold",
                 color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_fda_classes():
    fig, ax = plt.subplots(figsize=(11, 4.5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 5)
    classes = [
        (1.5, 2.5, "Class I\n(Low Risk)", "Documentation,\nformatting tools", M_GREEN, "Minimal\noversight"),
        (5.5, 2.5, "Class II\n(Moderate)", "Screening aids,\nmeasurement", M_ORANGE, "510(k)\nclearance"),
        (9.5, 2.5, "Class III\n(High Risk)", "Autonomous\ndiagnosis", M_RED, "Premarket\napproval (PMA)"),
    ]
    for x, y, name, example, color, pathway in classes:
        rect = mpatches.FancyBboxPatch((x - 1.5, y - 1.2), 3.0, 3.0, boxstyle="round,pad=0.2",
                                        facecolor=color, alpha=0.10, edgecolor=color, linewidth=2)
        ax.add_patch(rect)
        ax.text(x, y + 0.8, name, ha="center", fontsize=13, fontweight="bold", color=color)
        ax.text(x, y - 0.0, example, ha="center", fontsize=10, color=M_DARK)
        ax.text(x, y - 0.9, pathway, ha="center", fontsize=9, color=color, fontweight="bold")
    # arrow
    ax.annotate("", xy=(9.5, 0.8), xytext=(1.5, 0.8),
                arrowprops=dict(arrowstyle="-|>", lw=2.5, color=M_DARK))
    ax.text(5.5, 0.4, "Increasing Regulatory Burden  \u2192", ha="center",
            fontsize=11, color=M_DARK, fontweight="bold")
    ax.set_title("FDA Software as Medical Device (SaMD) Classification",
                 fontsize=14, fontweight="bold", color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

def fig_course_arc():
    fig, ax = plt.subplots(figsize=(12, 4.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 5)
    weeks = [
        (1.5, "Weeks 1-4", "Image\nProcessing", M_BLUE, "Pixels"),
        (4.5, "Weeks 5-10", "Deep\nLearning", M_GREEN, "Features"),
        (7.5, "Weeks 11-13", "Foundation\nModels", M_ORANGE, "Knowledge"),
        (10.5, "Weeks 14-16", "Agentic\nSystems", M_MAROON, "Orchestration"),
    ]
    for x, wk, name, color, concept in weeks:
        rect = mpatches.FancyBboxPatch((x - 1.2, 1.5), 2.4, 2.5, boxstyle="round,pad=0.2",
                                        facecolor=color, alpha=0.12, edgecolor=color, linewidth=2)
        ax.add_patch(rect)
        ax.text(x, 3.5, name, ha="center", fontsize=13, fontweight="bold", color=color)
        ax.text(x, 2.5, concept, ha="center", fontsize=11, color=M_DARK, style="italic")
        ax.text(x, 1.8, wk, ha="center", fontsize=9, color=M_GRAY)
    # arrows
    for i in range(3):
        x1 = weeks[i][0] + 1.3
        x2 = weeks[i+1][0] - 1.3
        ax.annotate("", xy=(x2, 2.75), xytext=(x1, 2.75),
                    arrowprops=dict(arrowstyle="-|>", lw=2, color=M_DARK))
    ax.text(6.0, 0.6, "The agent orchestrates everything from earlier themes",
            ha="center", fontsize=12, color=M_DARK, fontweight="bold", style="italic")
    ax.set_title("Course Arc: CSCI 3397", fontsize=14, fontweight="bold",
                 color=M_DARK, pad=10)
    ax.axis("off")
    fig.tight_layout()
    return fig

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE DECKS
# ══════════════════════════════════════════════════════════════════════════════

def build_lec32():
    prs = new_prs()
    # ── Title ──
    title_slide(prs, 32, "LLM APIs & the Agency Spectrum",
                "Theme 4: Agentic Systems \u2014 Building a Mini Claude Code (1/5)")
    # ── Housekeeping ──
    content_slide(prs, "Housekeeping", [
        "Recap: Lec. 30-31 \u2014 Multimodal foundation models (CLIP, LLaVA)",
        "Today: From using FMs to orchestrating them \u2014 giving LLMs agency",
        "Theme 4 plan: 5 lectures \u2192 build a mini Claude Code step-by-step",
        "    Reference implementation: nano-claude-code (~5K lines Python)",
        "Action: ps4 out (agentic systems)",
        "Reading: smolagents intro (HuggingFace); Anthropic API quickstart",
    ])
    # ── Section 1 ──
    section_slide(prs, "1", "The Agency Spectrum")
    figure_slide(prs, "Agency Is a Spectrum, Not Binary", fig_agency_spectrum(),
                 note="Source: adapted from HuggingFace smolagents conceptual guide")
    content_slide(prs, "The Multi-Step Agent Pattern", [
        "memory = [user_task]",
        "while llm_should_continue(memory):",
        "    action = llm_get_next_action(memory)",
        "    observations = execute_action(action)",
        "    memory += [action, observations]",
        "",
        "This is the core of what we're building \u2192 nano-claude-code's agent.py",
    ])
    content_slide(prs, "When to Use Agents vs. Simpler Approaches", [
        "\u2713  Workflow is predictable \u2192 hardcode it (100% reliable, no LLM errors)",
        "\u2717  User requests are unpredictable, require flexible tool combinations \u2192 agent",
        "",
        "Example: \"Is it possible to reschedule my MRI to Wednesday and still",
        "    get results before my Thursday appointment?\"",
        "    \u2192 No fixed workflow can handle this",
        "",
        "Biomedical motivation: clinical workflows involve many variable steps",
        "    (retrieve prior scans, compare, measure, report)",
    ])
    # ── Section 2 ──
    section_slide(prs, "2", "LLM API Basics")
    figure_slide(prs, "The Messages Format", fig_messages_format(),
                 note="Universal across Anthropic, OpenAI, Google \u2014 role-tagged message list")
    figure_slide(prs, "API Call Flow", fig_api_call_flow())
    code_slide(prs, "Making an API Call (Anthropic Python SDK)",
               'import anthropic\n'
               'client = anthropic.Anthropic()  # reads ANTHROPIC_API_KEY from env\n'
               '\n'
               'response = client.messages.create(\n'
               '    model="claude-sonnet-4-6-20250514",\n'
               '    max_tokens=1024,\n'
               '    system="You are a biomedical image analysis assistant.",\n'
               '    messages=[{\n'
               '        "role": "user",\n'
               '        "content": "What structures are visible in a brain MRI T1?"\n'
               '    }]\n'
               ')\n'
               'print(response.content[0].text)',
               note="nano-claude-code connection: this is providers.py \u2014 the API abstraction layer")
    code_slide(prs, "Streaming \u2014 Tokens Arrive Incrementally",
               'with client.messages.stream(\n'
               '    model="claude-sonnet-4-6-20250514",\n'
               '    max_tokens=1024,\n'
               '    messages=[{"role": "user", "content": "Describe this slide..."}]\n'
               ') as stream:\n'
               '    for text in stream.text_stream:\n'
               '        print(text, end="", flush=True)',
               note="Important for interactive UX \u2014 user sees output as it's generated")
    figure_slide(prs, "Key Parameters: Context Windows", fig_token_context(),
                 note="Tokens \u2248 len(text) / 3.5  |  Cost = input_tokens + output_tokens")
    # ── Section 3 ──
    section_slide(prs, "3", "Structured Outputs & Prompt Engineering")
    figure_slide(prs, "From Unstructured Reports to Structured JSON", fig_structured_output())
    content_slide(prs, "Prompt Engineering for Medical Domain", [
        "1. Be specific: \"Extract all anatomical structures\" > \"Analyze this report\"",
        "2. Provide examples (few-shot): show 1-2 input/output pairs",
        "3. Set constraints: \"Output only JSON. No explanations.\"",
        "4. Domain grounding: include medical context in system prompt",
        "5. Chain-of-thought: \"First identify findings, then classify severity, then format\"",
    ])
    code_slide(prs, "Vision API \u2014 Sending Images to the LLM",
               'response = client.messages.create(\n'
               '    model="claude-sonnet-4-6-20250514",\n'
               '    messages=[{\n'
               '        "role": "user",\n'
               '        "content": [\n'
               '            {"type": "image",\n'
               '             "source": {"type": "base64",\n'
               '                        "media_type": "image/png",\n'
               '                        "data": img_b64}},\n'
               '            {"type": "text",\n'
               '             "text": "Describe the histological features."}\n'
               '        ]\n'
               '    }]\n'
               ')',
               note="Direct connection to Lec. 30-31 (multimodal FMs)")
    # ── Wrap-up ──
    content_slide(prs, "Wrap-Up & Looking Ahead", [
        "Key takeaway: LLM APIs are the foundation \u2014 messages in, text/data out",
        "The agency spectrum shows how to give LLMs increasing control",
        "",
        "What we built today: the API layer (providers.py in nano-claude-code)",
        "",
        "Next (Lec. 33): Tool Use & the Core Agent Loop",
        "    \u2192 We teach the LLM to call functions, turning it into an agent",
    ])
    path = os.path.join(OUT_DIR, "l32_slides.pptx")
    prs.save(path)
    print(f"Saved {path}")

def build_lec33():
    prs = new_prs()
    title_slide(prs, 33, "Tool Use & the Core Agent Loop",
                "Theme 4: Agentic Systems \u2014 Building a Mini Claude Code (2/5)")
    content_slide(prs, "Housekeeping", [
        "Recap: Lec. 32 \u2014 LLM APIs, messages format, streaming, structured outputs",
        "Today: From \u2605\u2606\u2606 (text in/out) to \u2605\u2605\u2606 (LLM calls functions)",
        "    This is the core of what makes an agent",
        "Reading: Anthropic tool use docs",
        "Code: nano-claude-code agent.py (179 lines) + tool_registry.py (98 lines)",
    ])
    # ── Section 1 ──
    section_slide(prs, "1", "What Is Tool Use?")
    figure_slide(prs, "The Tool Use Cycle", fig_tool_use_cycle(),
                 note="The LLM doesn't execute anything \u2014 it only requests. Your code decides whether to run it.")
    code_slide(prs, "Defining a Tool Schema (JSON Schema)",
               'tools = [{\n'
               '    "name": "count_cells",\n'
               '    "description": "Count cells in a microscopy image\\n'\
               'using watershed segmentation.",\n'
               '    "input_schema": {\n'
               '        "type": "object",\n'
               '        "properties": {\n'
               '            "image_path": {\n'
               '                "type": "string",\n'
               '                "description": "Path to the microscopy image"\n'
               '            },\n'
               '            "threshold": {\n'
               '                "type": "number",\n'
               '                "description": "Detection threshold (0-1)"\n'
               '            }\n'
               '        },\n'
               '        "required": ["image_path"]\n'
               '    }\n'
               '}]',
               note="The LLM reads the description to decide when to use each tool")
    # ── Section 2 ──
    section_slide(prs, "2", "The Agent Loop \u2014 ReAct Pattern")
    figure_slide(prs, "ReAct: Reasoning + Acting (Yao et al., 2022)", fig_react_loop())
    code_slide(prs, "The Core Agent Loop (~30 lines of Python)",
               'def agent_loop(user_message, tools, system_prompt):\n'
               '    messages = [{"role": "user", "content": user_message}]\n'
               '\n'
               '    while True:\n'
               '        response = call_llm(messages, tools, system_prompt)\n'
               '        messages.append({"role": "assistant",\n'
               '                         "content": response.content})\n'
               '\n'
               '        if response.stop_reason != "tool_use":\n'
               '            return response.text  # Done!\n'
               '\n'
               '        for tool_call in response.tool_calls:\n'
               '            result = execute_tool(\n'
               '                tool_call.name, tool_call.input)\n'
               '            messages.append({\n'
               '                "role": "user",\n'
               '                "content": [{"type": "tool_result",\n'
               '                    "tool_use_id": tool_call.id,\n'
               '                    "content": result}]\n'
               '            })\n'
               '        # Loop: LLM sees results, decides next action',
               note="nano-claude-code's agent.py: 179 lines with streaming events (TextChunk, ToolStart, ToolEnd)")
    content_slide(prs, "Key Design Decisions in the Agent Loop", [
        "Loop termination: LLM decides when to stop (stop_reason != 'tool_use')",
        "    + max-iterations guard to prevent infinite loops",
        "",
        "Message accumulation: every action + observation appended",
        "    \u2192 LLM sees its full history (context window matters!)",
        "",
        "Multiple tool calls: LLM can request several tools in one turn",
        "    \u2192 parallel execution for speed",
        "",
        "Event-driven: agent.py uses Python generators (yield)",
        "    \u2192 decouples logic from UI (same loop for CLI, web, notebook)",
    ])
    # ── Section 3 ──
    section_slide(prs, "3", "The Tool Registry \u2014 Plugin Architecture")
    figure_slide(prs, "Tool Registry: Register Once, Use Everywhere", fig_tool_registry())
    code_slide(prs, "ToolDef: The Universal Tool Interface",
               '@dataclass\n'
               'class ToolDef:\n'
               '    name: str              # "count_cells"\n'
               '    schema: dict           # JSON schema sent to API\n'
               '    func: Callable         # (params, config) -> result str\n'
               '    read_only: bool        # True = safe, auto-approve\n'
               '    concurrent_safe: bool  # True = can run in parallel\n'
               '\n'
               '# Registry operations:\n'
               'register_tool(tool_def)       # add a new tool at runtime\n'
               'get_tool(name)                # look up by name\n'
               'get_tool_schemas()            # export all schemas for API\n'
               'execute_tool(name, params)    # dispatch + output truncation\n'
               '\n'
               '# Output truncation: if result > 32K chars,\n'
               '# keep first 50% + last 25%',
               note="tool_registry.py: 98 lines. Open/closed principle \u2014 add tools without modifying the loop")
    # ── Code vs JSON ──
    figure_slide(prs, "Code Agents vs. JSON Tool Calling", fig_code_vs_json(),
                 note="smolagents insight: LLMs are trained on code \u2014 writing Python is more natural than JSON")
    # ── Demo + Wrap-up ──
    content_slide(prs, "Demo: Agent with Biomedical Tools", [
        "3 tools: count_cells(), measure_area(), compute_statistics()",
        "",
        "User: \"Analyze slide_042.png \u2014 count cells, measure largest cluster, stats\"",
        "",
        "Agent execution:",
        "  Think \u2192 call count_cells(\"slide_042.png\")",
        "  Observe \u2192 \"Found 247 cells in 12 clusters\"",
        "  Think \u2192 call measure_area(\"slide_042.png\", label=3)",
        "  Observe \u2192 \"Cluster area: 1,340 \u03bcm\u00b2\"",
        "  Think \u2192 call compute_statistics([247, 1340, ...])",
        "  Observe \u2192 stats object",
        "  Done \u2192 final text report with summary",
    ])
    content_slide(prs, "Wrap-Up & Looking Ahead", [
        "Key takeaway: Tool use + agent loop = core of every coding agent",
        "    The LLM reasons, calls tools, observes results, and loops",
        "    Entire loop is ~30 lines of Python",
        "",
        "What we built: tool_registry.py + agent.py \u2014 the brain of nano-claude-code",
        "",
        "Next (Lec. 34): Real tools \u2014 file I/O (Read, Write, Edit),",
        "    Bash execution, system prompts, and context management",
    ])
    path = os.path.join(OUT_DIR, "l33_slides.pptx")
    prs.save(path)
    print(f"Saved {path}")

def build_lec34():
    prs = new_prs()
    title_slide(prs, 34, "Coding Agent with File I/O & Bash",
                "Theme 4: Agentic Systems \u2014 Building a Mini Claude Code (3/5)")
    content_slide(prs, "Housekeeping", [
        "Recap: Lec. 33 \u2014 Tool use, ReAct loop, tool registry, code vs JSON agents",
        "Today: Give the agent real power \u2014 file I/O + shell execution",
        "    \u2192 turns it into a coding agent that writes and runs its own scripts",
        "Code: nano-claude-code tools.py, context.py, compaction.py",
    ])
    # ── Section 1 ──
    section_slide(prs, "1", "File I/O Tools")
    figure_slide(prs, "The Tool Suite: Read, Write, Edit, Glob, Grep, Bash", fig_tool_suite())
    code_slide(prs, "Read Tool \u2014 Returns File Content with Line Numbers",
               '{"name": "Read",\n'
               ' "description": "Read a file. Returns content with line numbers.",\n'
               ' "input_schema": {\n'
               '     "properties": {\n'
               '         "file_path": {"type": "string"},\n'
               '         "offset":    {"type": "integer",\n'
               '                       "description": "Start line"},\n'
               '         "limit":     {"type": "integer",\n'
               '                       "description": "Num lines to read"}\n'
               '     },\n'
               '     "required": ["file_path"]\n'
               ' }}\n'
               '\n'
               '# Also reads: images (base64), PDFs, Jupyter notebooks\n'
               '# read_only=True \u2192 auto-approved in \"auto\" permission mode',
               note="Edit tool: exact string replacement (old_string \u2192 new_string). More surgical than Write.")
    # ── Section 2: Bash ──
    section_slide(prs, "2", "The Bash Tool & Permission Gates")
    content_slide(prs, "Bash Tool \u2014 Power and Danger", [
        "The agent can now:",
        "  \u2022 Run Python scripts it just wrote:  python analyze_cells.py",
        "  \u2022 Install packages:  pip install scikit-image",
        "  \u2022 Git commands:  git diff, git log",
        "  \u2022 Data processing:  ffmpeg, imagemagick, ls -la",
        "",
        "But also could:  rm -rf /  or  curl malicious_url | bash",
        "",
        "\u2192 This is why permission gates are essential",
    ])
    figure_slide(prs, "Three Permission Modes", fig_permission_model())
    # ── Section 3: System Prompt ──
    section_slide(prs, "3", "System Prompt Engineering")
    figure_slide(prs, "Building the System Prompt \u2014 context.py", fig_system_prompt())
    content_slide(prs, "Domain-Specific System Prompts for Biomedical AI", [
        "Include domain knowledge:",
        "    \"This project analyzes H&E stained histopathology slides.\"",
        "    \"Standard magnifications: 5x, 10x, 20x, 40x.\"",
        "",
        "Include safety constraints:",
        "    \"Never delete patient data.\"",
        "    \"Always confirm before overwriting analysis results.\"",
        "",
        "Include conventions:",
        "    \"Use scikit-image for segmentation, pandas for stats.\"",
        "",
        "Good tool descriptions are critical \u2014 the LLM reads them to decide",
        "    which tool to call",
    ])
    # ── Section 4: Compaction ──
    section_slide(prs, "4", "Context Window Management")
    figure_slide(prs, "Two-Layer Context Compaction", fig_compaction(),
                 note="compaction.py (196 lines). Without this, long analysis sessions crash.")
    # ── Wrap-up ──
    content_slide(prs, "Demo: Agent Writes & Runs Analysis Script", [
        "User: \"Read all .tiff files in slides/, write a Python script that applies",
        "    Otsu thresholding, counts cells, saves results to CSV, show stats.\"",
        "",
        "Agent execution:",
        "  1. Glob(\"slides/*.tiff\") \u2192 finds 10 slides",
        "  2. Read(sample file) \u2192 inspects format",
        "  3. Write(\"analyze.py\") \u2192 creates Python script",
        "  4. Bash(\"python analyze.py\") \u2192 runs it",
        "  5. \u274c Script error \u2192 Read(traceback) \u2192 Edit(\"analyze.py\") \u2192 re-run",
        "  6. Read(\"results.csv\") \u2192 final text summary",
        "",
        "The agent debugs its own code iteratively!",
    ])
    content_slide(prs, "Wrap-Up & Looking Ahead", [
        "What we built: tools.py + context.py + compaction.py",
        "    Our mini Claude Code can now read, write, search, and execute",
        "",
        "Key insight: coding agent = agent loop + tools + system prompt + compaction",
        "    These 4 components (< 700 lines) give 80% of Claude Code's capability",
        "",
        "Next (Lec. 35): RAG, persistent memory, and multi-agent systems",
        "    \u2192 search literature, remember past sessions, delegate to sub-agents",
    ])
    path = os.path.join(OUT_DIR, "l34_slides.pptx")
    prs.save(path)
    print(f"Saved {path}")

def build_lec35():
    prs = new_prs()
    title_slide(prs, 35, "RAG, Memory & Multi-Agent Systems",
                "Theme 4: Agentic Systems \u2014 Building a Mini Claude Code (4/5)")
    content_slide(prs, "Housekeeping", [
        "Recap: Lec. 34 \u2014 File I/O tools, Bash, system prompt, context compaction",
        "Today: Three capabilities for complex workflows:",
        "    1. RAG \u2014 retrieve external knowledge",
        "    2. Memory \u2014 persist across sessions",
        "    3. Multi-agent \u2014 delegate to specialized sub-agents",
        "Code: nano-claude-code memory/, multi_agent/, mcp/",
    ])
    # ── Section 1: RAG ──
    section_slide(prs, "1", "RAG \u2014 Retrieval-Augmented Generation")
    figure_slide(prs, "The RAG Pipeline", fig_rag_pipeline())
    code_slide(prs, "RAG as a Tool (Not a Standalone System)",
               '# RAG is just another tool the agent can call:\n'
               'tools = [{\n'
               '    "name": "SearchLiterature",\n'
               '    "description": "Search indexed medical literature.\\n'\
               'Returns relevant passages with citations.",\n'
               '    "input_schema": {\n'
               '        "properties": {\n'
               '            "query": {"type": "string"}\n'
               '        }\n'
               '    }\n'
               '}]\n'
               '\n'
               '# The agent decides WHEN to search:\n'
               '# 1. Read the slide\n'
               '# 2. Form a hypothesis\n'
               '# 3. Search for supporting literature\n'
               '# Much more flexible than a fixed RAG pipeline!',
               note="Biomedical RAG: PubMed abstracts, radiology archives, clinical protocols")
    content_slide(prs, "Biomedical RAG Examples", [
        "Index radiology archive \u2192 \"Find prior reports with similar findings\"",
        "Index PubMed abstracts \u2192 \"Differential diagnosis for 2cm hypodense liver lesion?\"",
        "Index clinical protocols \u2192 \"Standard staining protocol for Ki-67?\"",
        "",
        "Embedding models for biomedical text:",
        "    General: OpenAI text-embedding-3-small, Cohere embed-v3",
        "    Domain-specific: PubMedBERT embeddings, BioLORD",
        "    \u2192 Evaluate retrieval quality on your domain before committing",
    ])
    # ── Section 2: Memory ──
    section_slide(prs, "2", "Persistent Memory")
    figure_slide(prs, "Four Memory Types", fig_memory_types())
    content_slide(prs, "Memory in Practice", [
        "Each memory is a Markdown file with YAML frontmatter:",
        "    ---",
        "    name: staining_protocol",
        "    type: project",
        "    ---",
        "    Use 10-min hematoxylin, 3-min eosin. Deparaffinize at 60\u00b0C.",
        "",
        "Memory tools: MemorySave, MemorySearch, MemoryDelete",
        "",
        "User says: \"Remember we use 20x magnification for this cohort\"",
        "    \u2192 Agent saves it, applies in all future sessions",
        "",
        "Memory vs RAG:",
        "    Memory = small, curated, personal, always loaded",
        "    RAG = large external corpora, searched on demand",
    ])
    # ── Section 3: Multi-agent ──
    section_slide(prs, "3", "Multi-Agent Systems")
    figure_slide(prs, "Multi-Agent Pathology Pipeline", fig_multi_agent())
    code_slide(prs, "Sub-Agent Implementation (ThreadPoolExecutor)",
               'class SubAgentManager:\n'
               '    def __init__(self, max_workers=3):\n'
               '        self.executor = ThreadPoolExecutor(\n'
               '            max_workers=max_workers)\n'
               '\n'
               '    def spawn(self, prompt, system_prompt, tools):\n'
               '        future = self.executor.submit(\n'
               '            agent_loop, prompt, tools, system_prompt)\n'
               '        return task_id\n'
               '\n'
               '# Depth limiting: max 3 levels of nesting\n'
               '#   agent \u2192 sub-agent \u2192 sub-sub-agent (max)\n'
               '#   Prevents infinite recursion\n'
               '\n'
               '# Built-in agent types:\n'
               '#   researcher \u2014 read-only tools\n'
               '#   coder      \u2014 read + write tools\n'
               '#   reviewer   \u2014 code review with severity levels',
               note="nano-claude-code multi_agent/: cooperative cancellation, fresh context per sub-agent")
    # ── MCP ──
    figure_slide(prs, "MCP: Connecting Agents to Clinical Systems", fig_mcp_protocol(),
                 note="Standard protocol \u2014 auditable access without embedding credentials in the agent")
    # ── Wrap-up ──
    content_slide(prs, "Wrap-Up & Looking Ahead", [
        "RAG \u2192 external knowledge    Memory \u2192 session continuity",
        "Sub-agents \u2192 parallelism     MCP \u2192 external systems",
        "",
        "Together: single-loop agent \u2192 system for real clinical workflows",
        "",
        "When NOT to use multi-agent: if the task is sequential and simple,",
        "    one agent is better \u2014 sub-agents add latency + complexity",
        "",
        "Our mini Claude Code now has all major components!",
        "",
        "Next (Lec. 36): Full architecture walkthrough, deployment",
        "    (HIPAA, FDA), and live demo of the complete system",
    ])
    path = os.path.join(OUT_DIR, "l35_slides.pptx")
    prs.save(path)
    print(f"Saved {path}")

def build_lec36():
    prs = new_prs()
    title_slide(prs, 36, "Putting It All Together\nAgentic Biomedical AI",
                "Theme 4: Agentic Systems \u2014 Building a Mini Claude Code (5/5)")
    content_slide(prs, "Housekeeping", [
        "What we built over Lec. 32-35:",
        "    32: LLM APIs & agency spectrum (providers.py)",
        "    33: Tool use & agent loop (tool_registry.py, agent.py)",
        "    34: File I/O, Bash, system prompts, compaction (tools.py, context.py, compaction.py)",
        "    35: RAG, memory, multi-agent, MCP (memory/, multi_agent/, mcp/)",
        "",
        "Today: Full system, agent landscape, deployment, live demo",
        "Action: ps4 due; final project slides due",
    ])
    # ── Section 1 ──
    section_slide(prs, "1", "The Full Architecture")
    figure_slide(prs, "nano-claude-code: Complete Architecture", fig_full_architecture())
    content_slide(prs, "REPL, Skills & Plugins", [
        "REPL (nano_claude.py \u2014 1,985 lines):",
        "    Interactive loop: read input \u2192 run agent \u2192 stream output",
        "    Slash commands: /model, /config, /memory, /save, /tasks",
        "    Permission prompts: ask user before writes/execution",
        "",
        "Skills \u2014 reusable prompt templates (Markdown + YAML):",
        "    /commit \u2192 loads skill prompt \u2192 runs through agent loop",
        "    /analyze-slide \u2192 standard cell counting pipeline",
        "",
        "Plugins \u2014 runtime-extensible tools:",
        "    Install from git: dicom-tools adds ReadDICOM, ExtractMetadata",
    ])
    # ── Section 2 ──
    section_slide(prs, "2", "The Agent Landscape")
    figure_slide(prs, "Where Does Our Agent Fit?", fig_agent_landscape())
    content_slide(prs, "Architectural Approaches Compared", [
        "JSON tool calling (OpenAI, LangChain):",
        "    LLM outputs JSON \u2192 parse \u2192 execute. Safe, easy to validate.",
        "",
        "Code agents (smolagents, Claude Code):",
        "    LLM writes executable Python. More flexible, needs sandboxing.",
        "",
        "Graph-based (LangGraph):",
        "    Predefined state machine with LLM at decision nodes.",
        "    More predictable, less flexible.",
        "",
        "The trade-off: more agency = more capability but less predictability",
        "    For clinical use, find the right balance",
    ])
    # ── Section 3 ──
    section_slide(prs, "3", "Deploying Biomedical AI Agents")
    figure_slide(prs, "The Deployment Stack", fig_deployment_stack())
    figure_slide(prs, "FDA: Software as Medical Device (SaMD)", fig_fda_classes())
    content_slide(prs, "HIPAA & Patient Privacy", [
        "PHI: patient names, dates, MRNs, images with metadata",
        "",
        "Cloud API risk: sending PHI to Claude/GPT may violate HIPAA",
        "    unless a BAA (Business Associate Agreement) is in place",
        "",
        "Solutions:",
        "    \u2022 Use providers with BAA (Anthropic, Azure OpenAI, AWS Bedrock)",
        "    \u2022 De-identify data before API calls",
        "    \u2022 Run local models (Ollama, vLLM) on hospital infrastructure",
        "",
        "Our architecture supports this: providers.py abstracts over cloud & local",
        "    Switch Claude \u2192 Ollama with one config change",
    ])
    content_slide(prs, "Latency, Cost & Monitoring", [
        "Latency: agent loops = multiple round trips (5-step = 5 API calls)",
        "    \u2192 Use fast models for simple steps, parallel tool calls, caching",
        "",
        "Cost: typically $0.01-0.10 per analysis",
        "",
        "Monitoring:",
        "    \u2022 Log every tool call, LLM response, user interaction",
        "    \u2022 Track: token usage, latency, error rates, approval/rejection",
        "    \u2022 Drift detection: compare outputs against validation set",
        "",
        "Permission model in production:",
        "    auto: routine analyses   |   manual: anything touching records",
        "    Never accept-all in production",
    ])
    # ── Section 4 ──
    section_slide(prs, "4", "Live Demo")
    content_slide(prs, "End-to-End: Kidney Biopsy Cohort Analysis", [
        "1. Agent reads MEMORY.md \u2192 recalls project context, Otsu thresholding, 20x",
        "2. Glob(\"slides/*.tiff\") \u2192 finds 10 slides",
        "3. Writes count_glomeruli.py (scikit-image watershed)",
        "4. Bash(\"python count_glomeruli.py\") \u2192 counts and areas",
        "5. SearchLiterature(\"normal glomerular area\") \u2192 retrieves 3 papers",
        "6. Writes final Markdown report:",
        "    \u2022 Per-slide counts and areas",
        "    \u2022 Comparison to literature normals",
        "    \u2022 Flagged abnormal slides",
        "7. MemorySave \u2192 records analysis parameters for future sessions",
    ])
    # ── Wrap-up ──
    figure_slide(prs, "The Arc of CSCI 3397", fig_course_arc())
    content_slide(prs, "What We Built in 5 Lectures", [
        "Lec 32: providers.py (618 L) \u2014 talk to LLMs",
        "Lec 33: agent.py + tool_registry.py (277 L) \u2014 core loop + dispatch",
        "Lec 34: tools.py + context.py + compaction.py (1,425 L) \u2014 real tools",
        "Lec 35: memory/ + multi_agent/ + mcp/ (~3,000 L) \u2014 extensions",
        "Lec 36: full system (~5,000 L) \u2014 complete coding agent",
        "",
        "An AI agent is not magic \u2014",
        "    it's a while loop, a tool registry, and a good system prompt.",
        "",
        "Understanding the architecture demystifies tools like Claude Code",
        "    and empowers you to build domain-specific agents",
        "    for any biomedical workflow.",
    ])
    content_slide(prs, "Good Luck with Final Projects!", [
        "Presentations: Fri Apr 24, Mon Apr 27, Wed Apr 29",
        "",
        "The agent we built can USE everything from earlier themes:",
        "    \u2022 Image processing (weeks 1-4) \u2014 filtering, thresholding",
        "    \u2022 Deep learning (weeks 5-10) \u2014 U-Net, detection, segmentation",
        "    \u2022 Foundation models (weeks 11-13) \u2014 SAM, CLIP, LLaVA",
        "    \u2022 Agentic systems (weeks 14-16) \u2014 orchestrate all of the above",
        "",
        "Thank you for a great semester!",
    ])
    path = os.path.join(OUT_DIR, "l36_slides.pptx")
    prs.save(path)
    print(f"Saved {path}")

if __name__ == "__main__":
    build_lec32()
    build_lec33()
    build_lec34()
    build_lec35()
    build_lec36()
    print("\nAll 5 slide decks generated!")
